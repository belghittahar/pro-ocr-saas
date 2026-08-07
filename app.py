import streamlit as st
from streamlit_option_menu import option_menu
from streamlit_mic_recorder import speech_to_text
import io
import os
from PIL import Image
import docx
import pandas as pd
import google.generativeai as genai
import PyPDF2
import pdfplumber
from pptx import Presentation

# -----------------------------------------------------------------------------
# Configuration & State
# -----------------------------------------------------------------------------
st.set_page_config(
    page_title="Pixel2Word - AI Document Platform",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded"
)

if "show_camera" not in st.session_state:
    st.session_state.show_camera = False

# -----------------------------------------------------------------------------
# Clean Native Streamlit Layout 
# -----------------------------------------------------------------------------
st.markdown("""
<style>
/* Hide default footer, but KEEP header/hamburger menu */
footer {visibility: hidden;}

/* Minimal spacing adjustments */
.main-title {
    text-align: center;
    font-size: 2.5rem !important;
    font-weight: 800;
    margin-bottom: 0.5rem;
}

.sub-title {
    text-align: center;
    font-size: 1.1rem;
    margin-bottom: 3rem;
    opacity: 0.8;
}

/* Footer */
.premium-footer {
    text-align: center;
    padding: 2rem;
    font-size: 0.85rem;
    margin-top: 4rem;
    opacity: 0.7;
}

/* Subpages Container */
.page-container {
    max-width: 800px;
    margin: 0 auto;
    padding: 2rem;
}

/* Specifically target the 'Browse files' button inside the dropzone */
[data-testid="stFileUploadDropzone"] button {
    background-color: #ffffff !important;
    color: #111827 !important;
    border: 1px solid #d1d5db !important;
    border-radius: 6px !important;
    padding: 0.5rem 1rem !important;
    font-weight: 600 !important;
}
[data-testid="stFileUploadDropzone"] button:hover {
    background-color: #f3f4f6 !important;
    border-color: #9ca3af !important;
}
</style>
""", unsafe_allow_html=True)

# -----------------------------------------------------------------------------
# Core Extraction Functions
# -----------------------------------------------------------------------------
def extract_text_from_pdf(file_bytes):
    text = ""
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
    except Exception:
        try:
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        except Exception as e2:
            st.error(f"Failed to extract PDF: {str(e2)}")
    return text

def extract_text_from_docx(file_bytes):
    text = ""
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        st.error(f"Failed to extract Word document: {str(e)}")
    return text

def extract_text_from_excel(file_bytes):
    try:
        df = pd.read_excel(io.BytesIO(file_bytes))
        return df.to_string()
    except Exception as e:
        st.error(f"Failed to extract Excel document: {str(e)}")
        return ""

def extract_text_from_csv(file_bytes):
    try:
        df = pd.read_csv(io.BytesIO(file_bytes))
        return df.to_string()
    except Exception as e:
        st.error(f"Failed to extract CSV document: {str(e)}")
        return ""

def extract_text_from_pptx(file_bytes):
    text = ""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        st.error(f"Failed to extract PowerPoint document: {str(e)}")
    return text

def analyze_document_with_ai(file_content, is_image, api_key):
    genai.configure(api_key=api_key)
    
    prompt = """
    You are a strict, highly precise B2B Document Intelligence API. Your ONLY objective is to extract and structure data exactly as it appears in the document.

    STRICT RULES:
    1. ZERO CONVERSATIONAL FILLER: Return ONLY the extracted text and structured data. Do not include greetings, explanations, or concluding remarks like "Here is the text".
    2. PRESERVE SPATIAL LAYOUT & TABLES: You must meticulously maintain the original formatting, alignment, and line breaks. If you detect a table, invoice, or structured list, you MUST output it as a cleanly formatted Markdown table or perfectly spaced text that visually mirrors the original image/document structure.
    3. STRICT FIDELITY: Transcribe exactly what is written. Do NOT hallucinate, guess, infer, or add any information that is not explicitly visible in the document.

    REQUIRED OUTPUT FORMAT:
    ### Document Type
    (State the type, e.g., Invoice, Contract, Receipt, Spreadsheet)

    ### Language
    (State the auto-detected language)

    ### Key Entities
    (List critical entities found, e.g., Total Amount, Date, Invoice Number)

    ### Full Extracted Text
    (Provide the exact transcribed text and markdown tables here, maintaining structural layout)
    """

    try:
        model = genai.GenerativeModel('gemini-3.5-flash')
        if is_image:
            image = Image.open(io.BytesIO(file_content))
            response = model.generate_content([prompt, image])
        else:
            full_prompt = f"{prompt}\n\nDocument Text:\n{file_content}"
            response = model.generate_content(full_prompt)
        return response.text
    except Exception as e:
        st.error(f"AI Processing Error: {str(e)}")
        return None

def create_word_doc(text):
    doc = docx.Document()
    doc.add_paragraph(text)
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

def create_excel_doc(text):
    lines = text.split('\n')
    data = [[line] for line in lines if line.strip()]
    df = pd.DataFrame(data, columns=["Extracted Content"])
    
    bio = io.BytesIO()
    with pd.ExcelWriter(bio, engine='openpyxl') as writer:
        df.to_excel(writer, index=False)
    return bio.getvalue()

# -----------------------------------------------------------------------------
# Modular Pages
# -----------------------------------------------------------------------------
def render_image_to_text(gemini_api_key):
    st.markdown('<p class="main-title">Image to Text (OCR)</p>', unsafe_allow_html=True)
    st.markdown('<p class="sub-title">Extract structured data instantly from physical receipts and images.</p>', unsafe_allow_html=True)
    
    if not gemini_api_key:
        st.warning("⚠️ Gemini API Key is missing. Please set the 'GEMINI_API_KEY' environment variable.")
        st.stop()

    st.markdown("### Upload Image")
    
    uploaded_file = None
    
    col_upload, col_spacer, col_camera = st.columns([1, 0.1, 1])
    with col_upload:
        upload_input = st.file_uploader("Upload Image (JPG, PNG, WEBP)", type=["jpg", "jpeg", "png", "webp"], label_visibility="collapsed")
        if upload_input:
            uploaded_file = upload_input
            st.session_state.show_camera = False

    with col_camera:
        st.write("") 
        st.write("")
        if st.button("📷 Take Photo", key="btn_camera"):
            st.session_state.show_camera = not st.session_state.show_camera
            
        if st.session_state.show_camera:
            st.markdown("<br>", unsafe_allow_html=True)
            camera_input = st.camera_input("Capture Document", label_visibility="collapsed")
            if camera_input:
                uploaded_file = camera_input
                
    st.markdown("---")

    if uploaded_file is not None:
        st.markdown("### Analysis Results")
        
        file_bytes = uploaded_file.getvalue()
        
        res_col1, res_col2 = st.columns([1, 1.5], gap="medium")
        
        with res_col1:
            try:
                image = Image.open(uploaded_file)
                st.image(image, caption="Source Image", use_container_width=True, clamp=True)
            except:
                st.warning("Preview not available.")
            
        with res_col2:
            with st.spinner("Processing image..."):
                ai_extracted_data = analyze_document_with_ai(file_bytes, True, gemini_api_key)
                
            if ai_extracted_data:
                st.success("Analysis Complete")
                edited_text = st.text_area("Extracted Intelligence", value=ai_extracted_data, height=400, label_visibility="collapsed")
                
                # --- Voice-to-Text Feature ---
                st.markdown("<p style='font-size: 0.9rem; margin-top: 0.5rem;'>🎤 Click to dictate additional notes:</p>", unsafe_allow_html=True)
                dictated_text = speech_to_text(language='en', use_container_width=True, just_once=True, key='STT_img')
                if dictated_text:
                    edited_text = edited_text + "\n\n### Dictated Notes\n" + dictated_text
                    st.success("Notes appended!")

                st.markdown("<br><b>Export Data</b>", unsafe_allow_html=True)
                btn_col1, btn_col2 = st.columns(2)
                
                with btn_col1:
                    word_bytes = create_word_doc(edited_text)
                    st.download_button("📄 Download Word", data=word_bytes, file_name="pixel2word_ocr.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
                with btn_col2:
                    excel_bytes = create_excel_doc(edited_text)
                    st.download_button("📊 Download Excel", data=excel_bytes, file_name="pixel2word_ocr.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

def render_document_to_text(gemini_api_key):
    st.markdown('<p class="main-title">Document to Text</p>', unsafe_allow_html=True)
    st.markdown('<p class="sub-title">Intelligently extract data from PDFs, Word docs, and Excel sheets.</p>', unsafe_allow_html=True)
    
    if not gemini_api_key:
        st.warning("⚠️ Gemini API Key is missing. Please set the 'GEMINI_API_KEY' environment variable.")
        st.stop()

    st.markdown("### Upload Document")
    uploaded_file = st.file_uploader("Upload Document (PDF, DOCX, XLSX, CSV, PPTX, TXT)", type=["pdf", "docx", "xlsx", "csv", "txt", "pptx"])

    if uploaded_file is not None:
        st.markdown("---")
        st.markdown("### Analysis Results")
        
        file_bytes = uploaded_file.getvalue()
        file_extension = uploaded_file.name.split('.')[-1].lower()
        
        st.info(f"📄 Document Loaded: {uploaded_file.name}")
        
        with st.spinner("Parsing and analyzing document..."):
            extracted_raw_text = ""
            if file_extension == 'pdf':
                extracted_raw_text = extract_text_from_pdf(file_bytes)
            elif file_extension == 'docx':
                extracted_raw_text = extract_text_from_docx(file_bytes)
            elif file_extension == 'xlsx':
                extracted_raw_text = extract_text_from_excel(file_bytes)
            elif file_extension == 'csv':
                extracted_raw_text = extract_text_from_csv(file_bytes)
            elif file_extension == 'pptx':
                extracted_raw_text = extract_text_from_pptx(file_bytes)
            elif file_extension == 'txt':
                extracted_raw_text = file_bytes.decode('utf-8', errors='ignore')
            
            if not extracted_raw_text.strip():
                st.error("Failed to extract raw text from document, or document is empty.")
                ai_extracted_data = None
            else:
                ai_extracted_data = analyze_document_with_ai(extracted_raw_text, False, gemini_api_key)
            
        if ai_extracted_data:
            st.success("Analysis Complete")
            edited_text = st.text_area("Extracted Intelligence", value=ai_extracted_data, height=400, label_visibility="collapsed")
            
            # --- Voice-to-Text Feature ---
            st.markdown("<p style='font-size: 0.9rem; margin-top: 0.5rem;'>🎤 Click to dictate additional notes:</p>", unsafe_allow_html=True)
            dictated_text = speech_to_text(language='en', use_container_width=True, just_once=True, key='STT_doc')
            if dictated_text:
                edited_text = edited_text + "\n\n### Dictated Notes\n" + dictated_text
                st.success("Notes appended!")

            st.markdown("<br><b>Export Data</b>", unsafe_allow_html=True)
            btn_col1, btn_col2 = st.columns(2)
            
            with btn_col1:
                word_bytes = create_word_doc(edited_text)
                st.download_button("📄 Download Word", data=word_bytes, file_name="pixel2word_doc.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            with btn_col2:
                excel_bytes = create_excel_doc(edited_text)
                st.download_button("📊 Download Excel", data=excel_bytes, file_name="pixel2word_doc.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

def render_audio_to_text():
    st.markdown('<p class="main-title">Audio to Text</p>', unsafe_allow_html=True)
    st.markdown('<p class="sub-title">Dictate notes instantly using your microphone.</p>', unsafe_allow_html=True)

    st.markdown("### 🎙️ Voice Dictation")
    st.markdown("Click the button below to start recording. Speak clearly into your microphone.")
    
    # We place the mic recorder cleanly in the center
    dictated_text = speech_to_text(language='en', use_container_width=True, just_once=True, key='STT_main')
    
    if dictated_text:
        st.markdown("---")
        st.success("Audio transcribed successfully!")
        edited_text = st.text_area("Transcription", value=dictated_text, height=200, label_visibility="collapsed")
        
        st.markdown("<br><b>Export Data</b>", unsafe_allow_html=True)
        btn_col1, btn_col2 = st.columns(2)
        with btn_col1:
            word_bytes = create_word_doc(edited_text)
            st.download_button("📄 Download Word", data=word_bytes, file_name="pixel2word_audio.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        with btn_col2:
            st.download_button("📝 Download Text", data=edited_text, file_name="pixel2word_audio.txt", mime="text/plain")

def render_privacy():
    # Removed whitespace indentation to prevent rendering as code blocks
    st.markdown("""
<div class="page-container">
<h2>Privacy Policy</h2>
<p><em>Effective Date: 2024</em></p>
<p>Welcome to <strong>www.pixel2word.com</strong>. Enterprise data security is our highest priority.</p>

<h3>1. Data Processing and Security</h3>
<p>All files, images, and documents uploaded to our platform are processed securely in real-time. 
<strong>We DO NOT store, log, or train models on your documents.</strong> Following extraction, files are immediately discarded from active memory.</p>

<h3>2. Third-Party Infrastructure</h3>
<p>Our extraction pipeline relies on secure, enterprise-grade APIs (Google Generative AI). Data transmitted is subject to stringent enterprise security agreements ensuring absolute confidentiality.</p>

<h3>3. Analytics</h3>
<p>We utilize standard web analytics strictly to monitor platform health. This data remains anonymized and is never linked to the contents of uploaded documents.</p>
</div>
""", unsafe_allow_html=True)

def render_terms():
    # Removed whitespace indentation to prevent rendering as code blocks
    st.markdown("""
<div class="page-container">
<h2>Terms of Service</h2>
<p>By accessing <strong>www.pixel2word.com</strong>, you agree to comply with these Terms of Service.</p>

<h3>1. Use of Service</h3>
<p>Pixel2Word provides an advanced AI-powered document extraction tool provided "as is". While we leverage state-of-the-art AI to ensure accuracy, we do not guarantee 100% perfection. Users must verify extracted data prior to professional reliance.</p>

<h3>2. Acceptable Use</h3>
<p>You agree not to use the platform to upload illegal, explicit, or malicious files. We reserve the right to permanently revoke access for policy violations.</p>

<h3>3. Liability Limitation</h3>
<p>Pixel2Word and its operators shall not be held liable for direct or indirect damages resulting from data inaccuracies, service outages, or infrastructure limitations.</p>
</div>
""", unsafe_allow_html=True)

def render_contact():
    # Removed whitespace indentation to prevent rendering as code blocks
    st.markdown("""
<div class="page-container">
<h2>Contact Support</h2>
<p>Have technical questions, require API access, or need enterprise deployment support?</p>
<p>Reach out directly to our B2B integration team.</p>

<h3>Email Support</h3>
<p>📧 <strong>contact@pixel2word.com</strong></p>

<p><em>Standard SLA: Responses within 24 business hours.</em></p>
</div>
""", unsafe_allow_html=True)

# -----------------------------------------------------------------------------
# Main Application Router
# -----------------------------------------------------------------------------
def main():
    gemini_api_key = os.environ.get("GEMINI_API_KEY")

    with st.sidebar:
        st.markdown("<h2 style='text-align: center; margin-bottom: 2rem;'>Pixel2Word</h2>", unsafe_allow_html=True)
        
        # Premium Navigation using option_menu
        page = option_menu(
            menu_title=None,
            options=["Image to Text", "Document to Text", "Audio to Text", "Privacy Policy", "Terms of Service", "Contact Us"],
            icons=["camera", "file-earmark-text", "mic", "shield-lock", "file-earmark-check", "envelope"],
            menu_icon="cast",
            default_index=0,
            styles={
                "container": {"padding": "0!important", "background-color": "transparent", "border": "none"},
                "icon": {"font-size": "1.1rem"},
                "nav-link": {
                    "font-size": "0.95rem", 
                    "text-align": "left", 
                    "margin": "0.2rem 0", 
                    "border-radius": "6px",
                    "font-weight": "500"
                },
                "nav-link-selected": {
                    "font-weight": "600"
                },
            }
        )

    # Route Rendering
    if page == "Image to Text":
        render_image_to_text(gemini_api_key)
    elif page == "Document to Text":
        render_document_to_text(gemini_api_key)
    elif page == "Audio to Text":
        render_audio_to_text()
    elif page == "Privacy Policy":
        render_privacy()
    elif page == "Terms of Service":
        render_terms()
    elif page == "Contact Us":
        render_contact()

    # Premium Global Footer
    st.markdown("""
    <div class="premium-footer">
        <strong>© 2024 Pixel2Word. All rights reserved.</strong><br>
        Enterprise Document Intelligence
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
